import json
import os
import uuid
from typing import Sequence, Optional, Any

import fitz
import pandas as pd
from paddleocr import PaddleOCR
from pptx import Presentation
from tqdm import tqdm

from core.custom_types import DocumentMetadata
from llama_index.core.schema import Document as LlamaDocument
from docx import Document
import re

from core.settings import settings


def load_txt(path,metadata: DocumentMetadata)->Sequence[LlamaDocument]:
    """获取普通文本Document"""
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()
    metadata.section_title = ".".join(metadata.file_name.split(".")[:-1])
    return [
        LlamaDocument(
            text=text,
            metadata=metadata.dict()
        )
    ]

def load_docx(file_path: str, metadata: DocumentMetadata)->Sequence[LlamaDocument]:
    doc = Document(file_path)
    sections = []
    current_section = ""
    current_title = None
    for para in doc.paragraphs:
        text = para.text.strip()

        if not text:
            continue

        # 如果是标题
        if para.style.name.startswith("Heading"):

            if current_section:
                sections.append((current_title,current_section))

            current_title = text
            current_section = f"# {text}\n"
        else:
            current_section += text + "\n"

    if current_section:
        sections.append((current_title,current_section))

    # 解析表格
    for table in doc.tables:

        table_text = []

        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]

            table_text.append(" | ".join(row_data))

        sections.append((current_title,"\n".join(table_text)))

    documents = []

    for title,section in sections:
        metadata.section_title = title
        documents.append(
            LlamaDocument(
                text=section,
                metadata=metadata.dict()
            )
        )

    return documents



def extract_md_title(text: str) -> str | None:
    match = re.search(r'^(#{1,6})\s+(.*)', text, re.MULTILINE)
    if match:
        return match.group(2).strip()
    return None

def load_markdown(path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    """获取markdown Document"""
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()

    sections = []
    # 按 markdown 标题切分
    parts = re.split(r'(?=\n#{1,6} )', text)

    metadata.source = metadata.file_type
    for part in parts:
        part = part.strip()
        if not part:
            continue
        metadata.section_title = extract_md_title(part)
        sections.append(
            LlamaDocument(
                text=part,
                metadata=metadata.dict()
            )
        )

    return sections



ocr = PaddleOCR()

def extract_pdf_title(text: str) -> bool:
    """
    简单标题识别
    """
    if len(text) < 40 and re.match(r"^[0-9一二三四五六七八九十\.、 ]+", text):
        return True

    if len(text) < 20 and text.isupper():
        return True

    return False

def ocr_page(page):
    pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))

    img_path = f"temp_{uuid.uuid4().hex}.png"
    pix.save(img_path)

    result = ocr.ocr(img_path)
    texts = []
    for line in result[0]:
        bbox, (text, score) = line
        # if score > 0.3:
            # 去掉极短、无效字符
        text = text.replace("\n", " ").strip()
        if len(text) > 1:
            texts.append(text)
    try:
        os.remove(img_path)
    except:
        pass
    return "\n".join(texts)



def load_pdf(path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    """获取pdf Document"""
    pdf = fitz.open(path)
    documents = []

    current_section = ""
    current_title = ""
    source = metadata.file_type
    for page_num, page in tqdm(enumerate(pdf),desc="加载pdf"):

        blocks  = page.get_text("blocks")
        source =  metadata.file_type
        if len(blocks)<=1 or  sum(len(b[4]) for b in blocks) < 50:
            print(f"第{page_num+1}页扫描")
            #进行ocr扫描图文
            text = ocr_page(page)
            text = text.replace("\n", " ")
            blocks = [(0, 0, 0, 0, text, 0, 0)]
            source = "ocr"
        for block in blocks:
            text = block[4].strip()
            if not text:
                continue
            text = text.replace("\n", " ")

            # 标题、页码、页脚
            if len(text) < 5:
                continue

            if extract_pdf_title(text):
                if current_section:
                    documents.append(
                        LlamaDocument(
                            text=current_section,
                            metadata={
                                **metadata.dict(),
                                "page": page_num + 1,
                                "section_title": current_title,
                                "source":source
                            }
                        )
                    )
                current_title = text
                current_section = f"# {text}\n"
            else:
                current_section += text + "\n"

    if current_section:
        metadata.section_title =  current_title
        metadata.source =  metadata.file_type
        documents.append(
            LlamaDocument(
                text=current_section,
                metadata=metadata.dict()
            )
        )
    return documents


def load_excel(file_path: str, metadata: DocumentMetadata,header_mode=settings.excel_header_mode) -> Sequence[LlamaDocument]:
    """
    将 Excel 文件解析为 LlamaDocument
    每个 sheet 的每一行生成一个 document
    header_mode = "row" | 'col' | None
    """
    documents = []
    xls = pd.ExcelFile(file_path)

    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name, dtype=str)  # 全部读取为字符串
        df.fillna("", inplace=True)  # 空值处理
        header = ""
        data = df.iterrows()
        row_text  = ""
        documents = []
        for row_idx, row in data:
            # 拼接每一行的文本
            row_text += " | ".join([str(cell) for cell in row])

            if not row_text.strip():
                continue  # 跳过空行

            if row_idx == 0 and header_mode is not None:
                if header_mode == "row":
                    header = row_text
                row_text  = ""
                continue

            if len(row_text) >= settings.excel_min_chunk_size:
                # 生成 metadata
                doc_meta = {
                    **metadata.dict(),
                    "section_title": header,
                    "sheet_name":sheet_name,
                }
                documents.append(
                    LlamaDocument(
                        text=row_text,
                        metadata=doc_meta
                    )
                )
                row_text = ""
            else:
                row_text+='\n'

    return documents



def load_pptx(file_path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    """
        获取pptx的document
    """
    prs = Presentation(file_path)
    documents = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if not slide_texts:
            continue


        # 将整个 slide 文本按行拼接或按段落拼接
        slide_content = "\n".join(slide_texts)
        # 过滤短文本
        if len(slide_content) < 30:
            continue
        # 过滤致谢页
        if any(k in slide_content.lower() for k in settings.pptx_filter_slider) and slide_idx == len(prs.slides) -1:
            continue

        doc_meta = {
            **metadata.dict(),
            "page": slide_idx + 1,
            "section_title": slide.shapes.title.text if slide.shapes.title else "",
            "source": "pptx"
        }
        documents.append(LlamaDocument(text=slide_content, metadata=doc_meta))
    return documents





def load_json(path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)

    nodes = []

    if isinstance(data,list):
        i = 0
        j = 1
        while i < len(data):
            while  i+j < len(nodes) and len(str(nodes[i:i+j])) < settings.json_min_chunk_size:
                j+=1
            nodes.append(
                LlamaDocument(text="\n".join(str(nodes[i:i+j])), metadata=metadata.dict())
            )
            i = i+j
            j = 1
    elif isinstance(data,dict):
        obj = {}
        for k,v in data.items():
            obj[k] =v
            if  len(str(obj)) >= settings.json_min_chunk_size:
                nodes.append(
                    LlamaDocument(text=str(obj), metadata=metadata.dict())
                )
                obj = {}
        if not obj:
            nodes.append(
                LlamaDocument(text=str(obj), metadata=metadata.dict())
            )
    return nodes





def load_file(path:str, metadata: DocumentMetadata)->Sequence[LlamaDocument]:
    if path.endswith(".txt"):
        return load_txt(path,metadata)
    elif path.endswith(".docx") or path.endswith(".doc"):
        metadata.source = 'doc'
        return load_docx(path,metadata)
    elif path.endswith(".md") or path.endswith(".markdown"):
        metadata.source = 'markdown'
        return load_markdown(path,metadata)
    elif path.endswith(".pdf"):
        return load_pdf(path,metadata)
    elif path.endswith(".xlsx") or path.endswith(".xls") or path.endswith(".csv"):
        metadata.source = 'excel'
        return load_excel(path,metadata)
    elif path.endswith(".pptx") or  path.endswith(".ppt"):
        metadata.source = 'ppt'
        return load_pptx(path,metadata)
    elif path.endswith(".json"):
        return load_json(path,metadata)
    else:
        return []


