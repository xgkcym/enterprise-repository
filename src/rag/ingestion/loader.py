import json
import re
from collections import Counter
from typing import Sequence

import cv2
import fitz
import numpy as np
import pandas as pd
from docx import Document
from docx.document import Document as DocxDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from llama_index.core.schema import Document as LlamaDocument
from pptx import Presentation
from tqdm import tqdm

from core.custom_types import DocumentMetadata
from core.settings import settings
from src.rag.ocr_client import remote_ocr_image

def iter_docx_blocks(doc: DocxDocument):
    """Yield paragraphs and tables in their original document order."""
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield Table(child, doc)


def load_txt(path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    """Load a plain text file into a single document."""
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()
    metadata.source = metadata.source or metadata.file_type
    metadata.section_title = ".".join(metadata.file_name.split(".")[:-1])
    return [LlamaDocument(text=text, metadata=metadata.dict())]


def load_docx(file_path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    """Load a Word document by heading sections and tables."""
    doc = Document(file_path)
    sections = []
    current_section = ""
    current_title = None

    for block in iter_docx_blocks(doc):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text:
                continue

            if block.style and block.style.name.startswith("Heading"):
                if current_section:
                    sections.append((current_title, current_section))
                current_title = text
                current_section = f"# {text}\n"
            else:
                current_section += text + "\n"
            continue

        if isinstance(block, Table):
            table_rows = []
            for row in block.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                if any(cell for cell in row_data):
                    table_rows.append(" | ".join(row_data))

            table_text = "\n".join(table_rows).strip()
            if not table_text:
                continue

            if current_section:
                current_section += table_text + "\n"
            else:
                sections.append((current_title, table_text))

    if current_section:
        sections.append((current_title, current_section))

    documents = []
    for title, section in sections:
        metadata.section_title = title
        if len(section) < 10:
            continue
        documents.append(LlamaDocument(text=section, metadata=metadata.dict()))
    return documents


def extract_md_title(text: str) -> str | None:
    match = re.search(r'^(#{1,6})\s+(.*)', text, re.MULTILINE)
    if match:
        return match.group(2).strip()
    return None


def load_markdown(path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    """Load a markdown file by title sections."""
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()

    sections = []
    parts = re.split(r'(?=\n#{1,6} )', text)

    for part in parts:
        part = part.strip()
        if not part:
            continue
        metadata.section_title = extract_md_title(part)
        if len(part) < 10:
            continue
        sections.append(LlamaDocument(text=part, metadata=metadata.dict()))

    return sections


def ocr_image(img: np.ndarray, min_score: float | None = settings.orc_min_score) -> str:
    """Run OCR on an image and return extracted text."""
    score_threshold = 0.0 if min_score is None else min_score
    return remote_ocr_image(
        img,
        min_score=score_threshold,
        language=settings.orc_lang,
    )


def preprocess_image(img: np.ndarray) -> np.ndarray:
    """Ensure the image is converted to RGB/BGR-compatible format."""
    if img.ndim == 2:
        img = cv2.cvtColor(img, cv2.COLOR_GRAY2RGB)
    elif img.shape[2] == 4:
        img = cv2.cvtColor(img, cv2.COLOR_BGRA2BGR)
    return img


def extract_pdf_title(text: str) -> bool:
    """Use simple heuristics to detect PDF section titles."""
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return False
    # Bias toward citation correctness over aggressive title detection:
    # sentence-like fragments are treated as body text, not new sections.
    if re.search(r"[，,。；;：:！？!?]", text):
        return False
    if len(text) < 30 and re.match(r"^[0-9一二三四五六七八九十]+(?:\.[0-9]+)*[.、]?\s*\S+$", text):
        return True
    if len(text) < 15 and text.isupper() and re.search(r"[A-Z]", text):
        return True
    return False


def normalize_pdf_margin_text(text: str) -> str:
    """Normalize header/footer candidates for repeated-text detection."""
    normalized = re.sub(r"\s+", " ", text).strip()
    normalized = re.sub(r"\d+", "<num>", normalized)
    return normalized


def is_pdf_margin_block(block: tuple, page_height: float) -> bool:
    """Return True when a block sits in the top/bottom margin area."""
    y0, y1 = block[1], block[3]
    top_limit = page_height * 0.12
    bottom_limit = page_height * 0.88
    return y1 <= top_limit or y0 >= bottom_limit


def should_skip_pdf_margin_text(text: str, repeated_margin_texts: set[str]) -> bool:
    """Skip likely header/footer snippets and bare page-number text."""
    normalized = normalize_pdf_margin_text(text)
    if normalized in repeated_margin_texts:
        return True
    if re.fullmatch(r"(?:第\s*)?\d+\s*(?:页)?", text):
        return True
    if re.fullmatch(r"(?:page|p)\s*\d+(?:\s*(?:/|of)\s*\d+)?", text, re.IGNORECASE):
        return True
    return False


def load_pdf(path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    """Load a PDF file page by page, falling back to OCR when text extraction is weak."""
    pdf = fitz.open(path)
    documents = []
    current_title = ""
    pages_data = []
    margin_counter: Counter[str] = Counter()

    for page in pdf:
        blocks = page.get_text("blocks")
        source = metadata.file_type
        page_height = float(page.rect.height)
        if len(blocks) <= 1 or sum(len(block[4]) for block in blocks) < 50:
            pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
            img = preprocess_image(img)
            text = ocr_image(img)
            blocks = [(0, 0, 0, 0, text, 0, 0)]
            source = "ocr"
            page_height = float(pix.height)

        pages_data.append(
            {
                "blocks": blocks,
                "source": source,
                "page_height": page_height,
            }
        )

        for block in blocks:
            text = block[4].strip()
            if not text:
                continue
            text = text.replace("\n", " ").strip()
            if len(text) < 3 or len(text) > 120:
                continue
            if is_pdf_margin_block(block, page_height):
                margin_counter[normalize_pdf_margin_text(text)] += 1

    repeated_margin_texts = {text for text, count in margin_counter.items() if count >= 2}

    for page_num, page_data in tqdm(enumerate(pages_data), desc="加载 PDF"):
        blocks = page_data["blocks"]
        source = page_data["source"]
        page_height = page_data["page_height"]

        page_texts = []
        page_title = ""
        for block in blocks:
            text = block[4].strip()
            if not text:
                continue
            text = text.replace("\n", " ").strip()
            if len(text) < 5:
                continue
            if is_pdf_margin_block(block, page_height) and should_skip_pdf_margin_text(text, repeated_margin_texts):
                continue

            if not page_title and extract_pdf_title(text):
                page_title = text
                page_texts.append(f"# {text}")
            else:
                page_texts.append(text)

        page_content = "\n".join(page_texts).strip()
        if not page_content:
            continue

        if page_title:
            current_title = page_title

        documents.append(
            LlamaDocument(
                text=page_content,
                metadata={
                    **metadata.dict(),
                    "page": page_num + 1,
                    "section_title": page_title or current_title,
                    "source": source,
                },
            )
        )

    return documents


def load_excel(
    file_path: str,
    metadata: DocumentMetadata,
    header_mode=settings.excel_header_mode,
) -> Sequence[LlamaDocument]:
    """Load an Excel file and convert rows into Llama documents."""
    documents = []
    xls = pd.ExcelFile(file_path)

    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name, dtype=str)
        df.fillna("", inplace=True)
        header = ""
        row_text = ""

        for row_idx, row in df.iterrows():
            row_text += " | ".join([str(cell) for cell in row])
            if not row_text.strip():
                continue

            if row_idx == 0 and header_mode is not None:
                if header_mode == "row":
                    header = row_text
                row_text = ""
                continue

            if len(row_text) >= settings.excel_min_chunk_size:
                doc_meta = {
                    **metadata.dict(),
                    "section_title": header,
                    "sheet_name": sheet_name,
                }
                documents.append(LlamaDocument(text=row_text, metadata=doc_meta))
                row_text = ""
            else:
                row_text += "\n"

        if row_text.strip():
            doc_meta = {
                **metadata.dict(),
                "section_title": header,
                "sheet_name": sheet_name,
            }
            documents.append(LlamaDocument(text=row_text.strip(), metadata=doc_meta))

    return documents


def load_csv(
    file_path: str,
    metadata: DocumentMetadata,
    header_mode=settings.excel_header_mode,
) -> Sequence[LlamaDocument]:
    """Load a CSV file using the same chunking strategy as Excel."""
    documents = []
    df = pd.read_csv(file_path, dtype=str)
    df.fillna("", inplace=True)
    header = ""
    row_text = ""
    sheet_name = "csv"

    for row_idx, row in df.iterrows():
        row_text += " | ".join([str(cell) for cell in row])
        if not row_text.strip():
            continue

        if row_idx == 0 and header_mode is not None:
            if header_mode == "row":
                header = row_text
            row_text = ""
            continue

        if len(row_text) >= settings.excel_min_chunk_size:
            doc_meta = {
                **metadata.dict(),
                "section_title": header,
                "sheet_name": sheet_name,
            }
            documents.append(LlamaDocument(text=row_text, metadata=doc_meta))
            row_text = ""
        else:
            row_text += "\n"

    if row_text.strip():
        doc_meta = {
            **metadata.dict(),
            "section_title": header,
            "sheet_name": sheet_name,
        }
        documents.append(LlamaDocument(text=row_text.strip(), metadata=doc_meta))

    return documents


def load_pptx(file_path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    """Load a PowerPoint file by slide."""
    prs = Presentation(file_path)
    documents = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if not slide_texts:
            continue

        slide_content = "\n".join(slide_texts)
        if len(slide_content) < 30:
            continue
        if any(keyword in slide_content.lower() for keyword in settings.pptx_filter_slider) and slide_idx == len(prs.slides) - 1:
            continue

        doc_meta = {
            **metadata.dict(),
            "page": slide_idx + 1,
            "section_title": slide.shapes.title.text if slide.shapes.title else "",
            "source": "pptx",
        }
        documents.append(LlamaDocument(text=slide_content, metadata=doc_meta))
    return documents


def load_json(path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)

    nodes = []
    if isinstance(data, list):
        i = 0
        j = 1
        while i < len(data):
            while i + j < len(data) and len(json.dumps(data[i : i + j], ensure_ascii=False)) < settings.json_min_chunk_size:
                j += 1
            nodes.append(LlamaDocument(text=json.dumps(data[i : i + j], ensure_ascii=False), metadata=metadata.dict()))
            i = i + j
            j = 1
    elif isinstance(data, dict):
        obj = {}
        for key, value in data.items():
            obj[key] = value
            if len(json.dumps(obj, ensure_ascii=False)) >= settings.json_min_chunk_size:
                nodes.append(LlamaDocument(text=json.dumps(obj, ensure_ascii=False), metadata=metadata.dict()))
                obj = {}
        if obj:
            nodes.append(LlamaDocument(text=json.dumps(obj, ensure_ascii=False), metadata=metadata.dict()))
    return nodes


def load_image(path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    img = cv2.imread(path)
    if img is None:
        return []
    img = preprocess_image(img)
    text = ocr_image(img)
    if not text.strip():
        return []

    return [
        LlamaDocument(
            text=text,
            metadata={
                **metadata.dict(),
                "source": "ocr",
                "file_path": path,
            },
        )
    ]


def load_file(path: str, metadata: DocumentMetadata) -> Sequence[LlamaDocument]:
    if metadata.file_type in ["txt"]:
        return load_txt(path, metadata)
    if metadata.file_type in ["doc", "docx"]:
        metadata.source = "doc"
        return load_docx(path, metadata)
    if metadata.file_type in ["md", "markdown"]:
        metadata.source = "markdown"
        return load_markdown(path, metadata)
    if metadata.file_type in ["pdf"]:
        return load_pdf(path, metadata)
    if metadata.file_type in ["xls", "xlsx"]:
        metadata.source = "excel"
        return load_excel(path, metadata)
    if metadata.file_type in ["csv"]:
        metadata.source = "excel"
        return load_csv(path, metadata)
    if metadata.file_type in ["pptx"]:
        metadata.source = "pptx"
        return load_pptx(path, metadata)
    if metadata.file_type in ["json"]:
        metadata.source = "json"
        return load_json(path, metadata)
    if metadata.file_type in ["jpeg", "png", "jpg", "bmp", "webp", "tiff", "tif"]:
        return load_image(path, metadata)
    return []
