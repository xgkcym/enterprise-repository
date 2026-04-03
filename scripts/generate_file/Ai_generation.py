import os
import random
import json
from datetime import datetime, timedelta
import time

import pandas as pd
from faker import Faker
from reportlab.lib import colors
from tqdm import tqdm
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# AI改写相关
from openai import OpenAI
import dotenv
from tenacity import retry, stop_after_attempt, wait_exponential
dotenv.load_dotenv()
os.environ["PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"] = "True"

# ============================================
# 配置
# ============================================

fake = Faker("zh_CN")
ROOT = "data_ai"

# DeepSeek API配置（也可以用OpenAI）
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")
DEEPSEEK_BASE_URL = "https://api.deepseek.com"

# 目录结构保持不变
folders = [
    "01_company_profiles",
    "02_industry_research",
    "03_investment_memo",
    "04_fund_reports",
    "05_risk_management",
    "06_meeting_minutes",
    "07_portfolio_data",
    "08_market_data",
    "09_internal_training",
    "10_internal_emails"
]

for f in folders:
    os.makedirs(os.path.join(ROOT, f), exist_ok=True)


# ============================================
# AI内容生成器（核心替换）
# ============================================

class AIContentGenerator:
    """使用AI生成多样化的金融文本内容"""

    def __init__(self, api_type="deepseek"):
        self.api_type = api_type
        if api_type == "deepseek":
            self.client = OpenAI(
                api_key=DEEPSEEK_API_KEY,
                base_url=DEEPSEEK_BASE_URL
            )
        else:
            self.client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
    def generate_paragraph(self, context=None):
        """
        生成一个多样化的段落
        context: 可选的上下文信息，如公司名、行业等
        """
        if context is None:
            context = {}

        # 随机选择段落主题
        topics = [
            "公司介绍", "行业分析", "投资观点", "风险提示",
            "财务预测", "竞争格局", "政策影响", "技术趋势",
            "估值分析", "市场情绪", "供应链分析", "客户动态"
        ]
        topic = random.choice(topics)

        # 随机选择风格
        styles = ["专业严谨", "简洁明了", "数据驱动", "前瞻性", "保守稳健"]
        style = random.choice(styles)

        prompt = f"""请生成一段金融投资领域的专业文本，要求：

主题：{topic}
风格：{style}
长度：80-150字

{'额外信息：' + str(context) if context else ''}

要求：
1. 内容专业、数据可信
2. 如果是公司介绍，包含成立时间、员工数、营收等
3. 如果是行业分析，包含市场规模、增速、趋势等
4. 如果是投资观点，包含评级、目标价、逻辑等
5. 直接输出文本，不要任何解释和标记
6. 输出语言为中文

生成的文本："""

        try:
            response = self.client.chat.completions.create(
                model="deepseek-chat" if self.api_type == "deepseek" else "gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "你是一位资深的金融分析师，擅长撰写专业的投资研究报告。"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.8,  # 较高温度增加多样性
                max_tokens=300
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            print(f"AI生成失败: {e}, 使用备用文本")
            return self._fallback_paragraph()

    def _fallback_paragraph(self):
        """备用段落（API失败时使用）"""
        fallbacks = [
            "根据最新市场数据，该行业呈现出稳健增长态势，预计未来三年复合增长率将达到15-20%。",
            "公司核心业务保持强劲增长，市场份额持续提升，技术壁垒构筑了较强的竞争优势。",
            "宏观环境存在不确定性，但行业结构性机会依然明确，建议关注龙头企业。",
            "估值层面，当前市盈率处于历史中位数水平，具备一定的安全边际。",
        ]
        return random.choice(fallbacks)

    def generate_title(self, doc_type):
        """生成多样化的文档标题"""
        titles = {
            "md": ["行业深度研究报告", "赛道投资价值分析", "产业趋势展望", "细分领域专题研究"],
            "docx": ["投资研究报告", "尽职调查备忘录", "投资决策建议书", "项目评估报告"],
            "pdf": ["基金投资研究报告", "资产配置建议书", "季度投资回顾", "年度投资策略"],
            "ppt": ["投资策略分析", "行业机会解读", "组合管理汇报", "投决会材料"],
            "txt": ["内部通讯", "投资备忘录", "研究简报", "晨会纪要"]
        }
        return random.choice(titles.get(doc_type, ["研究报告"])) + f"_{datetime.now().strftime('%Y%m%d')}"

    def generate_section_title(self):
        """生成章节标题"""
        sections = [
            "市场概况", "竞争格局分析", "投资逻辑", "风险因素",
            "财务预测", "估值分析", "同业对比", "未来展望",
            "政策环境", "技术趋势", "供应链梳理", "客户分析"
        ]
        return random.choice(sections)

    def generate_conclusion(self):
        """生成投资结论"""
        conclusions = [
            "综合考虑行业景气度和公司竞争力，建议维持增持评级。",
            "短期波动不改长期趋势，建议逢低布局核心标的。",
            "行业处于成长期，龙头公司估值合理，建议重点关注。",
            "风险收益比较为均衡，建议保持标配配置。",
            "技术壁垒构筑护城河，看好公司长期发展空间。"
        ]
        return random.choice(conclusions)


# 全局AI生成器实例
ai_gen = AIContentGenerator("deepseek")

# ============================================
# 实体池（保留用于表格数据等）
# ============================================

industries = [
    "人工智能", "新能源", "半导体", "医疗科技", "机器人",
    "云计算", "消费科技", "区块链", "量子计算", "生物技术"
]

funds = [
    "HC Alpha Growth Fund", "HC AI Innovation Fund", "HC Global Macro Fund",
    "HC Energy Transition Fund", "HC Quant Arbitrage Fund", "HC Value Discovery Fund"
]

departments = [
    "投资策略部", "量化研究部", "风险管理部", "资产配置委员会", "投资委员会"
]

# 生成公司数据库（用于表格数据）
companies = []
for i in range(60):
    companies.append({
        "name": fake.company(),
        "industry": random.choice(industries),
        "founded": random.randint(1995, 2023),
        "employees": random.randint(50, 15000),
        "revenue": random.randint(5, 1000)
    })


# ============================================
# Markdown生成（使用AI内容）
# ============================================

def generate_md(path):
    with open(path, "w", encoding="utf8") as f:
        f.write(f"# {ai_gen.generate_title('md')}\n\n")
        f.write(f"**报告日期**: {datetime.now().strftime('%Y-%m-%d')}  \n")
        f.write(f"**分析师**: {fake.name()}  \n\n")
        f.write("---\n\n")

        num_sections = random.randint(4, 7)
        for i in range(num_sections):
            f.write(f"## {ai_gen.generate_section_title()}\n\n")

            # 每个章节3-6个段落
            for _ in range(random.randint(3, 6)):
                f.write(ai_gen.generate_paragraph() + "\n\n")

            # 添加表格
            f.write("### 重点公司对比\n\n")
            f.write("| 公司 | 行业 | 市值 | 增长率 |\n")
            f.write("|------|------|------|--------|\n")
            for _ in range(random.randint(8, 15)):
                c = random.choice(companies)
                f.write(f"| {c['name']} | {c['industry']} | {random.randint(50, 800)}亿 | {random.randint(5, 35)}% |\n")
            f.write("\n\n")

        f.write(f"## 投资结论\n\n")
        f.write(ai_gen.generate_conclusion() + "\n\n")


# ============================================
# DOCX生成（使用AI内容）
# ============================================

def generate_docx(path):
    doc = Document()

    doc.add_heading(ai_gen.generate_title('docx'), 0)
    doc.add_paragraph(f"报告日期：{datetime.now().strftime('%Y-%m-%d')}")
    doc.add_paragraph(f"分析师：{fake.name()}")
    doc.add_paragraph("")

    num_sections = random.randint(3, 6)
    for i in range(num_sections):
        doc.add_heading(ai_gen.generate_section_title(), level=1)

        for _ in range(random.randint(3, 5)):
            doc.add_paragraph(ai_gen.generate_paragraph())

        # 添加表格
        table = doc.add_table(rows=1, cols=4)
        table.style = 'Table Grid'
        hdr = table.rows[0].cells
        hdr[0].text = "公司"
        hdr[1].text = "行业"
        hdr[2].text = "市值"
        hdr[3].text = "增长率"

        for _ in range(random.randint(8, 12)):
            c = random.choice(companies)
            row = table.add_row().cells
            row[0].text = c['name']
            row[1].text = c['industry']
            row[2].text = str(random.randint(50, 800)) + "亿"
            row[3].text = str(random.randint(5, 35)) + "%"

    doc.add_heading("投资结论", level=1)
    doc.add_paragraph(ai_gen.generate_conclusion())

    doc.save(path)


# ============================================
# PDF生成（使用AI内容）
# ============================================

# 注册字体（保持原有逻辑）
FONT_PATH = "fonts/NotoSansCJKsc-VF.ttf"
if os.path.exists(FONT_PATH):
    pdfmetrics.registerFont(TTFont("NotoSans", FONT_PATH))
    style = ParagraphStyle("CN", fontName="NotoSans", fontSize=11, leading=16)
else:
    style = ParagraphStyle("CN", fontSize=11, leading=16)


def generate_pdf(path):
    from reportlab.platypus import PageBreak

    story = []

    # 标题
    story.append(Paragraph(ai_gen.generate_title('pdf'), style))
    story.append(Spacer(1, 20))
    story.append(Paragraph(f"报告日期：{datetime.now().strftime('%Y-%m-%d')}", style))
    story.append(Paragraph(f"分析师：{fake.name()}", style))
    story.append(Spacer(1, 20))

    # 正文段落
    for _ in range(random.randint(8, 15)):
        story.append(Paragraph(ai_gen.generate_paragraph(), style))
        story.append(Spacer(1, 10))

    # 表格
    table_data = [["公司", "行业", "市值", "增长率"]]
    for _ in range(random.randint(10, 15)):
        c = random.choice(companies)
        table_data.append([
            c['name'], c['industry'],
            str(random.randint(50, 800)) + "亿",
            str(random.randint(5, 35)) + "%"
        ])

    table = Table(table_data, colWidths=[6 * cm, 4 * cm, 3 * cm, 3 * cm])
    table.setStyle(TableStyle([
        ('FONTNAME', (0, 0), (-1, -1), 'NotoSans' if os.path.exists(FONT_PATH) else 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('ALIGN', (2, 1), (-1, -1), 'RIGHT'),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey)
    ]))
    story.append(table)
    story.append(Spacer(1, 20))

    # 结论
    story.append(Paragraph("投资结论", style))
    story.append(Spacer(1, 10))
    story.append(Paragraph(ai_gen.generate_conclusion(), style))

    pdf = SimpleDocTemplate(path)
    pdf.build(story)


# ============================================
# Excel生成（保持原有逻辑，表格数据）
# ============================================

def generate_excel(path):
    start_date = datetime(2024, 1, 1)
    months = 12
    data = []

    for c in companies[:30]:  # 限制公司数量避免文件过大
        company_name = c["name"]
        industry = c["industry"]

        prev_amount = random.randint(500000000, 5000000000)
        prev_return = round(random.uniform(-5, 5), 2)

        for month_offset in range(months):
            date = start_date + timedelta(days=30 * month_offset)
            date_str = date.strftime("%Y-%m-%d")

            change = random.uniform(-0.05, 0.05)
            amount = max(int(prev_amount * (1 + change)), 10000000)
            return_rate = round(prev_return + random.uniform(-2, 2), 2)

            data.append({
                "日期": date_str,
                "公司": company_name,
                "行业": industry,
                "持仓金额": amount,
                "收益率": return_rate
            })

            prev_amount = amount
            prev_return = return_rate

    df = pd.DataFrame(data)
    df.to_excel(path, index=False)


# ============================================
# JSON生成（保持原有逻辑）
# ============================================

def generate_json(path):
    data = []
    for _ in range(3000):
        c = random.choice(companies)
        data.append({
            "date": fake.date(),
            "company": c["name"],
            "industry": c["industry"],
            "price": round(random.uniform(20, 200), 2),
            "volume": random.randint(100000, 5000000)
        })

    with open(path, "w", encoding="utf8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


# ============================================
# PPT生成（使用AI内容）
# ============================================

def generate_ppt(path):
    prs = Presentation()

    # 标题页
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = ai_gen.generate_title('ppt')
    title_slide.placeholders[1].text = f"分析师：{fake.name()}\n日期：{datetime.now().strftime('%Y-%m-%d')}"

    # 内容页
    for i in range(random.randint(8, 12)):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = ai_gen.generate_section_title()

        content = slide.placeholders[1]
        paragraphs = [ai_gen.generate_paragraph() for _ in range(random.randint(2, 4))]
        content.text = "\n\n".join(paragraphs)

    prs.save(path)


# ============================================
# TXT生成（使用AI内容）
# ============================================

def generate_txt(path):
    with open(path, "w", encoding="utf8") as f:
        f.write(f"# {ai_gen.generate_title('txt')}\n\n")
        f.write(f"发件人：{fake.name()}\n")
        f.write(f"日期：{datetime.now().strftime('%Y-%m-%d')}\n")
        f.write("=" * 50 + "\n\n")

        for _ in range(random.randint(50, 100)):
            f.write(ai_gen.generate_paragraph())
            f.write("\n\n")


# ============================================
# 批量生成（保持原有调用方式）
# ============================================

def generate_all():
    print("开始生成AI增强的金融知识库...")
    print("=" * 50)

    print("\n[1/7] 生成Markdown报告...")
    for i in tqdm(range(1)):
        generate_md(f"{ROOT}/02_industry_research/industry_report_{i + 1}.md")

    print("\n[2/7] 生成DOCX备忘录...")
    for i in tqdm(range(1)):
        generate_docx(f"{ROOT}/03_investment_memo/investment_memo_{i + 1}.docx")

    print("\n[3/7] 生成PDF基金报告...")
    for i in tqdm(range(1)):
        generate_pdf(f"{ROOT}/04_fund_reports/fund_report_{i + 1}.pdf")

    print("\n[4/7] 生成Excel持仓数据...")
    for i in tqdm(range(1)):
        generate_excel(f"{ROOT}/07_portfolio_data/portfolio_data_{i + 1}.xlsx")

    print("\n[5/7] 生成JSON市场数据...")
    for i in tqdm(range(1)):
        generate_json(f"{ROOT}/08_market_data/market_data_{i + 1}.json")

    print("\n[6/7] 生成PPT培训材料...")
    for i in tqdm(range(1)):
        generate_ppt(f"{ROOT}/09_internal_training/training_{i + 1}.pptx")

    print("\n[7/7] 生成TXT内部邮件...")
    for i in tqdm(range(1)):
        generate_txt(f"{ROOT}/10_internal_emails/internal_mail_{i + 1}.txt")

    print("\n" + "=" * 50)
    print(f"✅ 金融RAG知识库生成完成！")
    print(f"📁 输出目录：{ROOT}")
    print(f"📊 共生成 {60 + 40 + 35 + 20 + 20 + 15 + 15} 个文件")


if __name__ == "__main__":
    generate_all()